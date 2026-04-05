from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLACK = RGBColor(0x0A, 0x0A, 0x0F)
DARK_BG = RGBColor(0x10, 0x10, 0x18)
CARD_BG = RGBColor(0x1A, 0x1A, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x88, 0x88, 0x99)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xFB, 0x71, 0x85)
SKY = RGBColor(0x38, 0xBD, 0xF8)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GRAY_LIGHT = RGBColor(0xD1, 0xD5, 0xDB)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
PINK = RGBColor(0xF4, 0x72, 0xB6)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
RED = RGBColor(0xF8, 0x71, 0x71)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=CARD_BG, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    else:
        shape.adjustments[0] = 0.02
    return shape


def add_accent_line(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT_GRAY, icon="✓", icon_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        # Icon
        run_icon = p.add_run()
        run_icon.text = f"{icon}  "
        run_icon.font.size = Pt(font_size)
        run_icon.font.color.rgb = icon_color or GOLD
        run_icon.font.name = 'Calibri'
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = 'Calibri'
    return txBox


# ============================================================
# SLIDE 1 — Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Gold accent line at top
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

# Decorative corner shape
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFB, 0xBF, 0x24)
shape.fill.fore_color.rgb.brightness = 0
shape.line.fill.background()

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
         "VISHVAKARMA HUB", 16, GOLD, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(1.2),
         "Vishvakarma Innovation\nChallenge 2026", 44, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
         "Sponsor Partnership Deck", 28, GOLD, False, PP_ALIGN.CENTER)

# Divider line
add_accent_line(slide, Inches(5.5), Inches(4.8), Inches(2.3), GOLD)

add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
         "India's Next Big Startup Movement — Be Part of It.", 18, MUTED, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
         "www.vishvakarmahub.com", 14, GOLD, False, PP_ALIGN.CENTER)

# Bottom accent
add_accent_line(slide, Inches(0), Inches(7.46), SLIDE_W, GOLD)


# ============================================================
# SLIDE 2 — The Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "THE OPPORTUNITY", 12, GOLD, True)

add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Why This Matters For Your Brand", 36, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
         "This is not just an event. It's a launchpad for India's next wave of innovation — and your brand can be at the center of it.",
         16, MUTED)

# Stat cards
stats = [
    ("10,000+", "Expected Attendees", "Students, founders, investors, industry leaders", AMBER),
    ("10,00,000+", "Digital Reach", "Social media impressions across platforms", SKY),
    ("500+", "Startups & Ideas", "Across AI, FinTech, HealthTech, EdTech & more", EMERALD),
    ("100+", "Colleges & Institutions", "Pan-India participation", VIOLET),
]

card_w = Inches(2.7)
card_h = Inches(2.2)
start_x = Inches(0.8)
card_gap = Inches(0.35)

for i, (val, title, sub, color) in enumerate(stats):
    x = start_x + i * (card_w + card_gap)
    y = Inches(2.6)
    card = add_shape_bg(slide, x, y, card_w, card_h, CARD_BG)
    add_accent_line(slide, x, y, card_w, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.35), card_w - Inches(0.6), Inches(0.6),
             val, 32, color, True, PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(0.4),
             title, 15, WHITE, True, PP_ALIGN.LEFT)
    add_text(slide, x + Inches(0.3), y + Inches(1.45), card_w - Inches(0.6), Inches(0.5),
             sub, 11, MUTED, False, PP_ALIGN.LEFT)

# Bottom highlights
highlights = [
    "🎤 Live Pitching Sessions",
    "🏪 Startup Exhibition",
    "🎓 Workshops & Masterclasses",
    "🤝 Investor Networking",
    "🏆 Grand Award Ceremony"
]
for i, h in enumerate(highlights):
    x = Inches(0.8) + i * Inches(2.4)
    add_text(slide, x, Inches(5.3), Inches(2.3), Inches(0.4), h, 13, LIGHT_GRAY, False)


# ============================================================
# SLIDE 3 — What Partners Get
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "PARTNERSHIP VALUE", 12, GOLD, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "What Every Partner Gets", 36, WHITE, True)

benefits = [
    ("🎯", "Brand Visibility", "Logo on stage, banners, website, certificates, ID cards, and all event media", AMBER),
    ("🎤", "Speaking Opportunities", "Keynote slots, panel discussions, and closing ceremony addresses", ROSE),
    ("💼", "Talent Access", "Direct access to 500+ innovative minds — recruit interns, co-founders, employees", EMERALD),
    ("📊", "Lead Generation", "QR-based lead capture, attendee database access, post-event analytics report", SKY),
    ("📱", "Digital Content", "Branded reels, YouTube coverage, social media posts reaching 10L+ impressions", PINK),
    ("🤝", "Networking", "VIP access to investor roundtables, founder meetups, exclusive after-event sessions", VIOLET),
]

for i, (emoji, title, desc, color) in enumerate(benefits):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.0)
    y = Inches(1.8) + row * Inches(2.5)
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(2.1), CARD_BG)
    add_accent_line(slide, x, y, Inches(3.7), color)
    add_text(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.6), Inches(0.5),
             emoji, 28, WHITE, False)
    add_text(slide, x + Inches(0.25), y + Inches(0.8), Inches(3.2), Inches(0.4),
             title, 18, WHITE, True)
    add_text(slide, x + Inches(0.25), y + Inches(1.25), Inches(3.2), Inches(0.7),
             desc, 12, MUTED)


# ============================================================
# SLIDE 4 — Title Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

add_text(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
         "MOST PREMIUM", 12, GOLD, True)

add_text(slide, Inches(0.8), Inches(0.8), Inches(6), Inches(0.8),
         "👑 Title Sponsor", 40, WHITE, True)

add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹7,50,000", 36, GOLD, True)

add_text(slide, Inches(0.8), Inches(2.15), Inches(8), Inches(0.5),
         '"The event carries YOUR name."', 16, MUTED, False)

benefits_left = [
    'Event named "Powered by [Your Brand]"',
    "Exclusive category rights — no competitors in your industry",
    "Logo on stage backdrop & all banners",
    "5–10 min keynote speech + closing ceremony address",
    "Dedicated hiring zone at venue",
]
benefits_right = [
    "Jury panel seat in finals",
    "Premium branding across website & all media",
    "Media coverage & press release mention",
    "Premium startup exhibition booth",
    "Direct access to top startups & talent pipeline",
]

add_shape_bg(slide, Inches(0.6), Inches(2.8), Inches(5.8), Inches(3.8), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.95), Inches(5.2), Inches(3.5),
                benefits_left, 14, LIGHT_GRAY, "✦", GOLD)

add_shape_bg(slide, Inches(6.7), Inches(2.8), Inches(5.8), Inches(3.8), CARD_BG)
add_bullet_list(slide, Inches(7.0), Inches(2.95), Inches(5.2), Inches(3.5),
                benefits_right, 14, LIGHT_GRAY, "✦", GOLD)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Own the event. Own the narrative.  →  Only 1 slot available.", 15, GOLD, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — Presenting Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, ROSE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
         "PREMIUM", 12, ROSE, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "🌟 Presenting Sponsor", 40, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹5,00,000", 36, ROSE, True)
add_text(slide, Inches(0.8), Inches(2.15), Inches(10), Inches(0.5),
         '"Co-host India\'s biggest innovation challenge."', 16, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.8), Inches(11.5), Inches(3.5), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.95), Inches(5.2), Inches(3.3), [
    'Co-host branding — full event co-presentation',
    'Sponsored challenge track (e.g., "AI Challenge powered by [You]")',
    'Logo on stage backdrop & event banners',
    '5 min keynote slot',
], 15, LIGHT_GRAY, "✦", ROSE)
add_bullet_list(slide, Inches(6.5), Inches(2.95), Inches(5.2), Inches(3.3), [
    'Premium branding on website & social media',
    'Media coverage & press mention',
    'VIP booth at startup exhibition',
    'Networking access with top founders',
], 15, LIGHT_GRAY, "✦", ROSE)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Be the co-host. Share the spotlight.", 15, ROSE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — Diamond Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, SKY)

add_text(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
         "ELITE", 12, SKY, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "💎 Diamond Sponsor", 40, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹3,50,000", 36, SKY, True)

add_shape_bg(slide, Inches(0.6), Inches(2.3), Inches(11.5), Inches(4.2), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8), [
    "Access to live startup pitching sessions",
    "Investor roundtable invite with top founders",
    "Lead capture system (QR code / digital cards)",
    "Logo on event banners and stage",
    "Featured website section with company profile",
], 15, LIGHT_GRAY, "✦", SKY)
add_bullet_list(slide, Inches(6.5), Inches(2.5), Inches(5.2), Inches(3.8), [
    "Social media promotion across all channels",
    "Premium exhibition booth",
    "VIP networking access",
    "Award ceremony mention & brand visibility",
], 15, LIGHT_GRAY, "✦", SKY)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Elite access. Maximum ROI.", 15, SKY, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — Platinum Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, PURPLE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
         "POPULAR CHOICE", 12, PURPLE, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "⬛ Platinum Sponsor", 40, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹1,00,000", 36, PURPLE, True)
add_text(slide, Inches(0.8), Inches(2.15), Inches(10), Inches(0.5),
         '"Maximum visibility at a smart investment."', 16, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.8), Inches(11.5), Inches(3.0), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.95), Inches(10), Inches(2.8), [
    "Logo on event banners",
    "Featured website placement",
    "Social media promotion across all platforms",
    "Booth at startup exhibition",
    "VIP networking access with founders & investors",
], 16, LIGHT_GRAY, "✦", PURPLE)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "The most popular tier. Smart brands choose Platinum.", 15, PURPLE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — Gold & Silver
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, ORANGE)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "VISIBILITY PACKAGES", 12, ORANGE, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Gold & Silver Sponsors", 36, WHITE, True)

# Gold card
add_shape_bg(slide, Inches(0.6), Inches(1.9), Inches(5.7), Inches(4.5), CARD_BG)
add_accent_line(slide, Inches(0.6), Inches(1.9), Inches(5.7), ORANGE)
add_text(slide, Inches(0.9), Inches(2.15), Inches(3), Inches(0.5),
         "🥇 Gold Sponsor", 24, WHITE, True)
add_text(slide, Inches(0.9), Inches(2.7), Inches(3), Inches(0.5),
         "₹50,000", 30, ORANGE, True)
add_bullet_list(slide, Inches(0.9), Inches(3.3), Inches(5.0), Inches(2.8), [
    "Logo on website",
    "Social media promotion",
    "Startup exhibition booth",
    "Event mention during ceremony",
], 15, LIGHT_GRAY, "✦", ORANGE)

# Silver card
add_shape_bg(slide, Inches(6.7), Inches(1.9), Inches(5.7), Inches(4.5), CARD_BG)
add_accent_line(slide, Inches(6.7), Inches(1.9), Inches(5.7), GRAY_LIGHT)
add_text(slide, Inches(7.0), Inches(2.15), Inches(3), Inches(0.5),
         "🥈 Silver Sponsor", 24, WHITE, True)
add_text(slide, Inches(7.0), Inches(2.7), Inches(3), Inches(0.5),
         "₹35,000", 30, GRAY_LIGHT, True)
add_bullet_list(slide, Inches(7.0), Inches(3.3), Inches(5.0), Inches(2.8), [
    "Logo on sponsor section",
    "Event promotion mention",
    "Networking access with founders",
], 15, LIGHT_GRAY, "✦", GRAY_LIGHT)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Great visibility. Budget-friendly.", 15, ORANGE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — Entry-Level Partners
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GREEN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "ENTRY-LEVEL PARTNERSHIPS", 12, GREEN, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Startup · Innovation · Community Partners", 32, WHITE, True)

# Startup Partner
add_shape_bg(slide, Inches(0.6), Inches(1.9), Inches(3.7), Inches(4.2), CARD_BG)
add_accent_line(slide, Inches(0.6), Inches(1.9), Inches(3.7), BLUE)
add_text(slide, Inches(0.9), Inches(2.15), Inches(3.2), Inches(0.5),
         "🚀 Startup Partner", 20, WHITE, True)
add_text(slide, Inches(0.9), Inches(2.65), Inches(3), Inches(0.5),
         "₹25,000", 28, BLUE, True)
add_bullet_list(slide, Inches(0.9), Inches(3.2), Inches(3.1), Inches(2.5), [
    "Logo on competition page",
    "Social media mention",
    "Access to startup database",
], 13, LIGHT_GRAY, "✦", BLUE)

# Innovation Partner
add_shape_bg(slide, Inches(4.7), Inches(1.9), Inches(3.7), Inches(4.2), CARD_BG)
add_accent_line(slide, Inches(4.7), Inches(1.9), Inches(3.7), GREEN)
add_text(slide, Inches(5.0), Inches(2.15), Inches(3.2), Inches(0.5),
         "💡 Innovation Partner", 20, WHITE, True)
add_text(slide, Inches(5.0), Inches(2.65), Inches(3), Inches(0.5),
         "₹15,000", 28, GREEN, True)
add_bullet_list(slide, Inches(5.0), Inches(3.2), Inches(3.1), Inches(2.5), [
    "Logo on event website",
    "Social media posts mention",
], 13, LIGHT_GRAY, "✦", GREEN)

# Community Partner
add_shape_bg(slide, Inches(8.8), Inches(1.9), Inches(3.7), Inches(4.2), CARD_BG)
add_accent_line(slide, Inches(8.8), Inches(1.9), Inches(3.7), PINK)
add_text(slide, Inches(9.1), Inches(2.15), Inches(3.2), Inches(0.5),
         "💜 Community Partner", 20, WHITE, True)
add_text(slide, Inches(9.1), Inches(2.65), Inches(3), Inches(0.5),
         "₹10,000", 28, PINK, True)
add_bullet_list(slide, Inches(9.1), Inches(3.2), Inches(3.1), Inches(2.5), [
    "Brand mention",
    "Website listing",
], 13, LIGHT_GRAY, "✦", PINK)

add_text(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.6),
         "Even at ₹10,000 — your brand reaches 10L+ people.\nThat's ₹0.01 per impression.", 14, MUTED, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — Innovation Track Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, VIOLET)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "STRATEGIC PARTNERSHIP", 12, VIOLET, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "🖥️ Innovation Track Sponsor", 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(4), Inches(0.6),
         "₹2,00,000 – ₹5,00,000", 30, VIOLET, True)
add_text(slide, Inches(0.8), Inches(2.1), Inches(10), Inches(0.4),
         "AI Track  ·  Robotics Track  ·  FinTech Track  ·  HealthTech Track", 14, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.7), Inches(11.5), Inches(3.5), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(3.2), [
    "Naming rights for your chosen track",
    "Direct access to niche startup talent",
    "Judging rights in track finals",
], 15, LIGHT_GRAY, "✦", VIOLET)
add_bullet_list(slide, Inches(6.5), Inches(2.9), Inches(5.2), Inches(3.2), [
    'Track winner announced as "[Your Brand] Award"',
    "Dedicated branding in track area",
    "Featured company profile on track page",
], 15, LIGHT_GRAY, "✦", VIOLET)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Own a track. Find your next billion-dollar idea before anyone else.", 15, VIOLET, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — Hiring Partner
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, EMERALD)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "HR BUDGETS", 12, EMERALD, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "💼 Hiring Partner", 40, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹3,00,000+", 36, EMERALD, True)
add_text(slide, Inches(0.8), Inches(2.15), Inches(10), Inches(0.5),
         '"Recruit from India\'s most innovative talent pool."', 16, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.8), Inches(11.5), Inches(3.5), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.95), Inches(5.2), Inches(3.2), [
    "Full resume database access of all participants",
    "On-spot interview booth at venue",
    "Branded as Official Hiring Partner",
], 15, LIGHT_GRAY, "✦", EMERALD)
add_bullet_list(slide, Inches(6.5), Inches(2.95), Inches(5.2), Inches(3.2), [
    "Job board placement on event website",
    "Priority access to winning teams",
    "Talent pipeline for internships & full-time roles",
], 15, LIGHT_GRAY, "✦", EMERALD)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "One event. Hundreds of pre-vetted innovators. Zero recruitment agency fees.", 15, EMERALD, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — Digital Reach Sponsor
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, PINK)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "MARKETING TEAMS", 12, PINK, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "📱 Digital Reach Sponsor", 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.6),
         "₹1,00,000 – ₹3,00,000", 30, PINK, True)
add_text(slide, Inches(0.8), Inches(2.1), Inches(10), Inches(0.5),
         '"Your brand in every reel, every story, every share."', 16, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.7), Inches(11.5), Inches(3.5), CARD_BG)
add_bullet_list(slide, Inches(0.9), Inches(2.9), Inches(5.2), Inches(3.2), [
    "Logo & branding in all Instagram reels & stories",
    "YouTube coverage with brand integration",
    "Influencer integration & co-created content",
], 15, LIGHT_GRAY, "✦", PINK)
add_bullet_list(slide, Inches(6.5), Inches(2.9), Inches(5.2), Inches(3.2), [
    "Branded hashtag campaign",
    "Post-event highlight reel with sponsor branding",
    "Social media analytics report shared post-event",
], 15, LIGHT_GRAY, "✦", PINK)

add_text(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Maximize your digital footprint. Measurable ROI.", 15, PINK, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — Special Partnerships
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, CYAN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "SPECIAL OPPORTUNITIES", 12, CYAN, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
         "Special Partnership Opportunities", 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
         "High visibility, low investment. Perfect add-ons for any brand.", 16, MUTED)

# Stage Sponsor
add_shape_bg(slide, Inches(0.6), Inches(2.3), Inches(3.7), Inches(4.0), CARD_BG)
add_accent_line(slide, Inches(0.6), Inches(2.3), Inches(3.7), CYAN)
add_text(slide, Inches(0.9), Inches(2.6), Inches(3.2), Inches(0.5),
         "🎤 Stage Sponsor", 22, WHITE, True)
add_text(slide, Inches(0.9), Inches(3.15), Inches(3), Inches(0.5),
         "₹40,000", 28, CYAN, True)
add_text(slide, Inches(0.9), Inches(3.75), Inches(3.1), Inches(1.2),
         "Your brand on the main stage backdrop. Every pitch, every talk, every photo — your logo is there.", 14, LIGHT_GRAY)

# Media Sponsor
add_shape_bg(slide, Inches(4.7), Inches(2.3), Inches(3.7), Inches(4.0), CARD_BG)
add_accent_line(slide, Inches(4.7), Inches(2.3), Inches(3.7), RED)
add_text(slide, Inches(5.0), Inches(2.6), Inches(3.2), Inches(0.5),
         "🎥 Media Sponsor", 22, WHITE, True)
add_text(slide, Inches(5.0), Inches(3.15), Inches(3), Inches(0.5),
         "₹30,000", 28, RED, True)
add_text(slide, Inches(5.0), Inches(3.75), Inches(3.1), Inches(1.2),
         "Logo in all videos and livestream. Your brand reaches everyone watching online.", 14, LIGHT_GRAY)

# Award Sponsor
add_shape_bg(slide, Inches(8.8), Inches(2.3), Inches(3.7), Inches(4.0), CARD_BG)
add_accent_line(slide, Inches(8.8), Inches(2.3), Inches(3.7), AMBER)
add_text(slide, Inches(9.1), Inches(2.6), Inches(3.2), Inches(0.5),
         "🏆 Award Sponsor", 22, WHITE, True)
add_text(slide, Inches(9.1), Inches(3.15), Inches(3), Inches(0.5),
         "₹20,000", 28, AMBER, True)
add_text(slide, Inches(9.1), Inches(3.75), Inches(3.1), Inches(1.2),
         "Your name on winner trophies. Associated with excellence and innovation forever.", 14, LIGHT_GRAY)


# ============================================================
# SLIDE 14 — Exhibition Booth
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, CYAN)

add_text(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
         "SHOWCASE YOUR PRODUCT", 12, CYAN, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
         "🏪 Exhibition Booth", 40, WHITE, True)
add_text(slide, Inches(0.8), Inches(1.55), Inches(4), Inches(0.6),
         "₹5,000", 36, CYAN, True)
add_text(slide, Inches(2.6), Inches(1.65), Inches(3), Inches(0.4),
         "per booth", 16, MUTED)

add_shape_bg(slide, Inches(0.6), Inches(2.4), Inches(11.5), Inches(3.5), CARD_BG)
add_accent_line(slide, Inches(0.6), Inches(2.4), Inches(11.5), CYAN)

features = [
    ("📐", "6×6 ft Branded Booth Space"),
    ("🪑", "Table, Chairs & Furniture"),
    ("⚡", "Power Outlet & Wi-Fi"),
    ("🎨", "Branded Backdrop"),
    ("👥", "10,000+ Visitor Footfall"),
    ("🎤", "Product Demo Space"),
]

for i, (emoji, feat) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.9) + col * Inches(3.8)
    y = Inches(2.7) + row * Inches(1.4)
    add_text(slide, x, y, Inches(0.5), Inches(0.4), emoji, 22, WHITE)
    add_text(slide, x + Inches(0.5), y + Inches(0.05), Inches(3.0), Inches(0.35),
             feat, 16, LIGHT_GRAY, True)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         "Showcase your product directly to 10,000+ attendees.\nPerfect for startups, product companies, and service providers.", 14, MUTED, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 — Tier Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

add_text(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.4),
         "AT A GLANCE", 12, GOLD, True)
add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.6),
         "Sponsorship Tier Comparison", 30, WHITE, True)

# Table header
headers = ["Benefit", "Title\n₹7.5L", "Presenting\n₹5L", "Diamond\n₹3.5L", "Platinum\n₹1L", "Gold\n₹50K", "Silver\n₹35K"]
col_widths = [Inches(2.6), Inches(1.5), Inches(1.6), Inches(1.6), Inches(1.5), Inches(1.4), Inches(1.4)]
start_x = Inches(0.6)
start_y = Inches(1.4)
row_h = Inches(0.48)

# Header row
x_pos = start_x
for j, (h, w) in enumerate(zip(headers, col_widths)):
    cell_bg = add_shape_bg(slide, x_pos, start_y, w, row_h, RGBColor(0x25, 0x25, 0x35))
    cell_bg.adjustments[0] = 0.0
    add_text(slide, x_pos + Inches(0.08), start_y + Inches(0.02), w - Inches(0.16), row_h,
             h, 9, GOLD, True, PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    x_pos += w

rows_data = [
    ("Event Naming", "✅", "—", "—", "—", "—", "—"),
    ("Keynote Speech", "10 min", "5 min", "—", "—", "—", "—"),
    ("Stage Logo", "✅", "✅", "✅", "—", "—", "—"),
    ("Banner Branding", "✅", "✅", "✅", "✅", "—", "—"),
    ("Exhibition Booth", "Premium", "VIP", "Premium", "Standard", "Standard", "—"),
    ("Website Branding", "Premium", "Featured", "Featured", "Featured", "Logo", "Logo"),
    ("Social Media", "All", "All", "All", "All", "Mention", "—"),
    ("Jury Seat", "✅", "—", "—", "—", "—", "—"),
    ("Hiring Zone", "✅", "—", "—", "—", "—", "—"),
    ("Talent Access", "✅", "✅", "✅", "✅", "—", "—"),
    ("Media Coverage", "✅", "✅", "—", "—", "—", "—"),
    ("Networking", "VIP", "VIP", "VIP", "VIP", "General", "General"),
]

for r, row in enumerate(rows_data):
    y = start_y + row_h + r * row_h
    x_pos = start_x
    bg_color = CARD_BG if r % 2 == 0 else RGBColor(0x14, 0x14, 0x1E)
    for j, (val, w) in enumerate(zip(row, col_widths)):
        cell_bg = add_shape_bg(slide, x_pos, y, w, row_h, bg_color)
        cell_bg.adjustments[0] = 0.0
        color = WHITE if j == 0 else (GREEN if val == "✅" else (MUTED if val == "—" else LIGHT_GRAY))
        add_text(slide, x_pos + Inches(0.08), y + Inches(0.05), w - Inches(0.16), row_h - Inches(0.1),
                 val, 9, color, j == 0, PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        x_pos += w


# ============================================================
# SLIDE 16 — Call to Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, GOLD)

add_text(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
         "Let's Build the\nFuture Together.", 48, WHITE, True, PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(2.6), Inches(2.3), GOLD)

add_text(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
         "India's innovators are waiting.\nYour brand can be the one that backed them.", 22, MUTED, False, PP_ALIGN.CENTER)

# Contact card
add_shape_bg(slide, Inches(3.5), Inches(4.2), Inches(6.3), Inches(2.0), CARD_BG)
add_accent_line(slide, Inches(3.5), Inches(4.2), Inches(6.3), GOLD)

add_text(slide, Inches(3.8), Inches(4.45), Inches(5.7), Inches(0.4),
         "Get In Touch", 20, GOLD, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3.8), Inches(4.9), Inches(5.7), Inches(0.35),
         "📧  sponsor@vishvakarmahub.com", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(3.8), Inches(5.25), Inches(5.7), Inches(0.35),
         "🌐  www.vishvakarmahub.com/competition/sponsor", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(3.8), Inches(5.6), Inches(5.7), Inches(0.35),
         "📞  Contact us for custom packages", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.5),
         "Limited slots  ·  First-come, first-served  ·  Secure your partnership today.", 14, GOLD, True, PP_ALIGN.CENTER)

add_accent_line(slide, Inches(0), Inches(7.46), SLIDE_W, GOLD)


# Save
output_path = r"c:\Users\psrih\Documents\VishvakarmaHub\Vishvakarma-Hub-Sponsor-Deck-2026.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
